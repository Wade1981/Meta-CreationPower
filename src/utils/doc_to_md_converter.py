#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
文档转换工具：将Word文档(.doc/.docx)和PowerPoint文档(.ppt)转换为Markdown格式

此工具可以在Enlightenment Lighthouse Runtime (ELR)中运行，
用于将Microsoft Office文档转换为Markdown格式，
便于在开源项目中使用和版本控制。
"""

import os
import sys
from pathlib import Path

# 尝试导入必要的库
try:
    from docx import Document
except ImportError:
    print("错误: 缺少python-docx库，请运行 'pip install python-docx' 安装")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("错误: 缺少python-pptx库，请运行 'pip install python-pptx' 安装")
    sys.exit(1)


def convert_docx_to_md(docx_path, md_path):
    """
    将Word文档(.docx)转换为Markdown格式
    
    Args:
        docx_path (str): Word文档路径
        md_path (str): 输出Markdown文件路径
    """
    try:
        # 打开Word文档
        doc = Document(docx_path)
        
        # 创建Markdown内容
        md_content = []
        
        # 遍历文档中的所有段落
        for para in doc.paragraphs:
            # 检查段落样式
            if para.style.name == 'Heading 1':
                md_content.append(f"# {para.text}\n")
            elif para.style.name == 'Heading 2':
                md_content.append(f"## {para.text}\n")
            elif para.style.name == 'Heading 3':
                md_content.append(f"### {para.text}\n")
            elif para.style.name == 'Heading 4':
                md_content.append(f"#### {para.text}\n")
            elif para.style.name == 'Heading 5':
                md_content.append(f"##### {para.text}\n")
            elif para.style.name == 'Heading 6':
                md_content.append(f"###### {para.text}\n")
            else:
                # 普通段落
                if para.text.strip():
                    md_content.append(f"{para.text}\n")
        
        # 遍历文档中的所有表格
        for table in doc.tables:
            # 开始表格
            md_content.append("| ")
            
            # 处理表头
            header_cells = table.rows[0].cells
            md_content.append(" | ".join([cell.text.strip() for cell in header_cells]))
            md_content.append(" |\n")
            
            # 表格分隔线
            md_content.append("| ")
            md_content.append(" | ".join(["---" for _ in header_cells]))
            md_content.append(" |\n")
            
            # 处理表格内容
            for row in table.rows[1:]:
                md_content.append("| ")
                md_content.append(" | ".join([cell.text.strip() for cell in row.cells]))
                md_content.append(" |\n")
            
            md_content.append("\n")
        
        # 将内容写入Markdown文件
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(''.join(md_content))
        
        print(f"成功: {docx_path} 已转换为 {md_path}")
        return True
        
    except Exception as e:
        print(f"错误: 转换Word文档时出错 - {str(e)}")
        return False


def convert_ppt_to_md(ppt_path, md_path):
    """
    将PowerPoint文档(.ppt)转换为Markdown格式
    
    Args:
        ppt_path (str): PowerPoint文档路径
        md_path (str): 输出Markdown文件路径
    """
    try:
        # 打开PowerPoint文档
        prs = Presentation(ppt_path)
        
        # 创建Markdown内容
        md_content = []
        
        # 遍历所有幻灯片
        for i, slide in enumerate(prs.slides, 1):
            # 添加幻灯片标题
            md_content.append(f"## 幻灯片 {i}\n")
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    # 添加形状文本
                    md_content.append(f"{shape.text_frame.text}\n")
            
            md_content.append("\n")
        
        # 将内容写入Markdown文件
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(''.join(md_content))
        
        print(f"成功: {ppt_path} 已转换为 {md_path}")
        return True
        
    except Exception as e:
        print(f"错误: 转换PowerPoint文档时出错 - {str(e)}")
        return False


def convert_to_md(input_path, output_path=None):
    """
    根据文件扩展名自动选择转换方法
    
    Args:
        input_path (str): 输入文件路径
        output_path (str): 输出Markdown文件路径，默认为输入文件路径替换扩展名
    """
    # 确保输入文件存在
    if not os.path.exists(input_path):
        print(f"错误: 输入文件不存在 - {input_path}")
        return False
    
    # 如果未指定输出路径，使用默认路径
    if output_path is None:
        base_path = os.path.splitext(input_path)[0]
        output_path = f"{base_path}.md"
    
    # 根据文件扩展名选择转换方法
    ext = os.path.splitext(input_path)[1].lower()
    
    if ext in ['.docx', '.doc']:
        return convert_docx_to_md(input_path, output_path)
    elif ext in ['.pptx', '.ppt']:
        return convert_ppt_to_md(input_path, output_path)
    else:
        print(f"错误: 不支持的文件格式 - {ext}")
        return False


def batch_convert(input_dir, output_dir=None):
    """
    批量转换目录中的所有支持的文档
    
    Args:
        input_dir (str): 输入目录路径
        output_dir (str): 输出目录路径，默认为输入目录
    """
    # 确保输入目录存在
    if not os.path.exists(input_dir):
        print(f"错误: 输入目录不存在 - {input_dir}")
        return False
    
    # 如果未指定输出目录，使用默认目录
    if output_dir is None:
        output_dir = input_dir
    else:
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
    
    # 支持的文件扩展名
    supported_extensions = ['.docx', '.doc', '.pptx', '.ppt']
    
    # 遍历目录中的所有文件
    for root, _, files in os.walk(input_dir):
        for file in files:
            # 检查文件扩展名
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_extensions:
                # 构建输入和输出路径
                input_path = os.path.join(root, file)
                rel_path = os.path.relpath(input_path, input_dir)
                output_file = f"{os.path.splitext(rel_path)[0]}.md"
                output_path = os.path.join(output_dir, output_file)
                
                # 确保输出目录存在
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                
                # 转换文件
                convert_to_md(input_path, output_path)
    
    print(f"批量转换完成: 已处理 {input_dir} 中的所有支持文件")
    return True


def main():
    """
    主函数，处理命令行参数
    """
    import argparse
    
    # 创建命令行参数解析器
    parser = argparse.ArgumentParser(description='将Word和PowerPoint文档转换为Markdown格式')
    
    # 添加命令行参数
    parser.add_argument('input', help='输入文件或目录路径')
    parser.add_argument('-o', '--output', help='输出文件或目录路径')
    parser.add_argument('-b', '--batch', action='store_true', help='批量转换目录中的所有文件')
    
    # 解析命令行参数
    args = parser.parse_args()
    
    # 根据参数执行转换
    if args.batch:
        # 批量转换
        batch_convert(args.input, args.output)
    else:
        # 单个文件转换
        convert_to_md(args.input, args.output)


if __name__ == "__main__":
    main()
